"""Office format loaders: XLSX (openpyxl) and PPTX (python-pptx)."""

from __future__ import annotations

import hashlib
import logging
from pathlib import Path

from docstack.models import DocumentRecord

logger = logging.getLogger(__name__)


def _doc_id(path: Path) -> str:
    st = path.stat()
    return hashlib.sha256(f"{path.resolve()}:{st.st_mtime_ns}".encode()).hexdigest()[:16]


def extract_xlsx_records(path: Path) -> list[DocumentRecord]:
    try:
        import openpyxl
    except ImportError as exc:
        raise RuntimeError("openpyxl is required for XLSX ingestion") from exc

    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    records: list[DocumentRecord] = []
    doc_id = _doc_id(path)

    for sheet_idx, sheet in enumerate(wb.worksheets, 1):
        rows: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c).strip() if c is not None else "" for c in row]
            line = "\t".join(cells).strip()
            if line:
                rows.append(line)
        if not rows:
            continue
        # 50-row blocks per sheet section
        block_size = 50
        for i in range(0, len(rows), block_size):
            text = "\n".join(rows[i : i + block_size])
            records.append(
                DocumentRecord(
                    doc_id=doc_id,
                    source_path=str(path.resolve()),
                    filename=path.name,
                    mime="xlsx",
                    page=sheet_idx,
                    block_type="table",
                    text=text,
                    section_heading=sheet.title,
                    table_id=f"sheet{sheet_idx}_block{i // block_size}",
                )
            )

    wb.close()
    return records


def extract_pptx_records(path: Path) -> list[DocumentRecord]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required for PPTX ingestion") from exc

    prs = Presentation(str(path))
    records: list[DocumentRecord] = []
    doc_id = _doc_id(path)

    for slide_num, slide in enumerate(prs.slides, 1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        parts.append(text)
            # Tables on slides
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    parts.append("\t".join(cells))

        text = "\n".join(parts).strip()
        if text:
            records.append(
                DocumentRecord(
                    doc_id=doc_id,
                    source_path=str(path.resolve()),
                    filename=path.name,
                    mime="pptx",
                    page=slide_num,
                    block_type="text",
                    text=text,
                )
            )

    return records
